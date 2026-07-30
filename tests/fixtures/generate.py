"""Generate test-fixture documents with real document libraries.

Run: python3 tests/fixtures/generate.py
Outputs land next to this file. Each fixture carries a unique phrase the
Rust extraction tests grep for; phrases must not repeat across fixtures.
"""

import json
from pathlib import Path

import docx
import openpyxl
import pptx
from odf.draw import Frame, Page, TextBox
from odf.opendocument import (
    OpenDocumentPresentation,
    OpenDocumentSpreadsheet,
    OpenDocumentText,
)
from odf.style import (
    DrawingPageProperties,
    MasterPage,
    PageLayout,
    PageLayoutProperties,
    Style,
)
from odf.table import Table, TableCell, TableRow
from odf.text import P

OUT = Path(__file__).parent


def make_docx():
    d = docx.Document()
    d.add_heading("Field Report", level=1)
    d.add_paragraph(
        "The team spent the morning walking the perimeter and noting which "
        "readings had drifted since the last visit."
    )
    d.add_paragraph(
        "Most of the afternoon went to quokka thermostat calibration, which "
        "took longer than planned but finished before the shift ended."
    )
    d.add_paragraph(
        "Remaining work is minor and can be folded into next week's rotation."
    )
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "Item"
    t.cell(0, 1).text = "Status"
    t.cell(1, 0).text = "ferret ledger entry"
    t.cell(1, 1).text = "Open"
    d.save(OUT / "notes.docx")


def make_pptx():
    p = pptx.Presentation()
    layout = p.slide_layouts[1]
    for i in range(1, 13):
        s = p.slides.add_slide(layout)
        s.shapes.title.text = f"Section {i}"
        body = s.placeholders[1].text_frame
        if i == 1:
            body.text = "opening remarks placeholder"
        elif i == 11:
            body.text = "marmoset budget review"
        else:
            body.text = f"Point {i}a covering the usual ground"
        body.add_paragraph().text = f"Follow-up item for section {i}"
    p.save(OUT / "deck.pptx")


def make_xlsx():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Summary"
    for row in [
        ["Region", "Units", "Revenue"],
        ["North", 128, 40960],
        ["South", 94, 30080],
        ["West", 211, 67520],
    ]:
        ws.append(row)
    q3 = wb.create_sheet("Q3 Notes")
    q3.append(["Note", "pangolin invoice discrepancy"])
    q3.append(["Amount", 1450.75])
    q3.append(["Days open", 12])
    wb.save(OUT / "sheet.xlsx")


def make_odt():
    d = OpenDocumentText()
    for text in [
        "Onboarding ran on schedule for the second quarter in a row.",
        "The tapir onboarding checklist was signed off by both reviewers.",
        "Nothing else is outstanding ahead of the quarterly close.",
    ]:
        d.text.addElement(P(text=text))
    d.save(OUT / "notes.odt", False)


def make_odp():
    d = OpenDocumentPresentation()
    layout = PageLayout(name="PL")
    layout.addElement(PageLayoutProperties(margin="0cm", pagewidth="28cm", pageheight="21cm"))
    d.automaticstyles.addElement(layout)
    master = MasterPage(name="Standard", pagelayoutname=layout)
    d.masterstyles.addElement(master)
    dp = Style(name="dp1", family="drawing-page")
    dp.addElement(DrawingPageProperties(backgroundsize="border"))
    d.automaticstyles.addElement(dp)

    for title, body in [
        ("Kickoff", "Scope agreed with both teams"),
        ("Timeline", "okapi rollout plan"),
        ("Risks", "Two dependencies still unconfirmed"),
    ]:
        page = Page(masterpagename=master, stylename=dp)
        d.presentation.addElement(page)
        for text, y in [(title, "2cm"), (body, "6cm")]:
            frame = Frame(width="24cm", height="3cm", x="2cm", y=y)
            page.addElement(frame)
            box = TextBox()
            frame.addElement(box)
            box.addElement(P(text=text))
    d.save(OUT / "deck.odp", False)


def make_ods():
    d = OpenDocumentSpreadsheet()
    table = Table(name="Expenses")
    for row in [
        ["Category", "Amount"],
        ["Travel", 820.5],
        ["civet expense summary", 1195.0],
        ["Supplies", 340.25],
    ]:
        tr = TableRow()
        table.addElement(tr)
        for value in row:
            if isinstance(value, str):
                cell = TableCell(valuetype="string")
            else:
                cell = TableCell(valuetype="float", value=value)
            cell.addElement(P(text=str(value)))
            tr.addElement(cell)
    d.spreadsheet.addElement(table)
    d.save(OUT / "sheet.ods", False)


def make_ipynb():
    nb = {
        "cells": [
            {
                "cell_type": "markdown",
                "metadata": {},
                "source": [
                    "# Sensor analysis\n",
                    "\n",
                    "This is the capybara regression writeup for the June batch.\n",
                ],
            },
            {
                "cell_type": "code",
                "execution_count": 1,
                "metadata": {"tags": ["calibration"]},
                "outputs": [
                    {
                        "name": "stdout",
                        "output_type": "stream",
                        "text": ["narwhal fit converged after 14 iterations\n"],
                    }
                ],
                "source": [
                    "def calibrate_axolotl(readings):\n",
                    "    # drop the warm-up samples before fitting\n",
                    "    trimmed = readings[5:]\n",
                    "    return sum(trimmed) / len(trimmed)\n",
                    "\n",
                    "print('narwhal fit converged after 14 iterations')\n",
                ],
            },
            {
                "cell_type": "markdown",
                "metadata": {},
                "source": [
                    "## Appendix\n",
                    "\n",
                    "See the quoll appendix notes for the raw sensor dumps.\n",
                ],
            },
        ],
        "metadata": {
            "kernelspec": {
                "display_name": "Python 3",
                "language": "python",
                "name": "python3",
            },
            "language_info": {
                "file_extension": ".py",
                "mimetype": "text/x-python",
                "name": "python",
                "nbconvert_exporter": "python",
                "pygments_lexer": "ipython3",
                "version": "3.11.4",
            },
        },
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    (OUT / "analysis.ipynb").write_text(json.dumps(nb, indent=1) + "\n")


if __name__ == "__main__":
    for fn in (make_docx, make_pptx, make_xlsx, make_odt, make_odp, make_ods, make_ipynb):
        fn()
        print(fn.__name__.removeprefix("make_"), "ok")
